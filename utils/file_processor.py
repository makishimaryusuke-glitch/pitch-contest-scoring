#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ファイル処理ユーティリティ
PDF、PowerPoint等のファイルからテキストを抽出します。
"""

import io
from pathlib import Path
from typing import Optional
import PyPDF2
from pptx import Presentation

def extract_text_from_pdf(file_path: Path) -> str:
    """PDFファイルからテキストを抽出"""
    try:
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"PDFのテキスト抽出に失敗しました: {e}")

def extract_text_from_pptx(file_path: Path) -> str:
    """PowerPointファイルからテキストを抽出"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"PowerPointのテキスト抽出に失敗しました: {e}")

def extract_text_from_file(file_path: Path) -> Optional[str]:
    """ファイル形式に応じてテキストを抽出"""
    suffix = file_path.suffix.lower()
    
    if suffix == '.pdf':
        return extract_text_from_pdf(file_path)
    elif suffix in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif suffix == '.txt':
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # UTF-8で読めない場合はShift-JISで試す
            with open(file_path, 'r', encoding='shift-jis') as f:
                return f.read()
    else:
        raise ValueError(f"サポートされていないファイル形式です: {suffix}")

def save_uploaded_file(uploaded_file, upload_dir: Path) -> Path:
    """アップロードされたファイルを保存"""
    upload_dir.mkdir(parents=True, exist_ok=True)
    file_path = upload_dir / uploaded_file.name
    
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    
    return file_path

def get_file_size(file_path: Path) -> int:
    """ファイルサイズを取得（バイト）"""
    return file_path.stat().st_size

def get_file_type(file_path: Path) -> str:
    """ファイルタイプを取得"""
    suffix = file_path.suffix.lower()
    type_map = {
        '.pdf': 'PDF',
        '.pptx': 'PowerPoint',
        '.ppt': 'PowerPoint',
        '.txt': 'Text',
        '.docx': 'Word',
        '.doc': 'Word',
    }
    return type_map.get(suffix, 'Unknown')















